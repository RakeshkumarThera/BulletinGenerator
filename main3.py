import os
import re
import subprocess
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from docx import Document
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# -------- CONFIG --------
CONFIG = {
    "SERVICE_ACCOUNT_FILE": "credentials.json",
    "LYRICS_FOLDER_NAME": "MPBC Song Lyric Docs",
    "TEMPLATE_PPTX": "BulletinTemplate.pptx",
    "OUTPUT_PPTX": "Updated_Bulletin.pptx",
    "FONT_NAME": "Montserrat",
    "FONT_SIZE": 12,
    "DATE_FONT_SIZE": 10,
    "TITLE_SIZE": 12,
    "FOOTNOTE_SIZE": 5,
    "BOX_TOP": 0.1,
    "BOX_WIDTH": 4.8,
    "BOX_HEIGHT": 7,
    "LEFT_MARGIN": 0.1,
    "RIGHT_MARGIN": 5.0,
    "DEFAULT_SONG_ORDER": [
        "Because He Lives",
        "A Mighty Fortress is Our God",
        "Christ the Lord is Risen Today",
        "Christ Arose",
        "See What a Morning (Resurrection Hymn)"
    ],
    "PAGE_NUMBERS": {
        (1, 'right'): 7,
        (2, 'left'): 6,
        (2, 'right'): 3,
        (3, 'left'): 4,
        (3, 'right'): 5
    }
}

# -------- GOOGLE SETUP --------
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
creds = service_account.Credentials.from_service_account_file(CONFIG["SERVICE_ACCOUNT_FILE"], scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

# -------- HELPERS --------
# def find_file_id(name, parent_id=None):
#     query = f"name = '{name}'"
#     if parent_id:
#         query += f" and '{parent_id}' in parents"
#     results = drive_service.files().list(q=query, fields="files(id, name)").execute()
#     return results.get("files", [])[0]["id"] if results.get("files") else None

def find_file_id(name, parent_id=None):
    # Escape apostrophes correctly for Google Drive query: use single backslash
    safe_name = name.replace("'", r"\'")
    query = f"name = '{safe_name}'"
    if parent_id:
        query += f" and '{parent_id}' in parents"
    try:
        print(f"[DEBUG] Querying Drive for: {safe_name}")
        results = drive_service.files().list(q=query, fields="files(id, name)").execute()
        files = results.get('files', [])
        return files[0]['id'] if files else None
    except Exception as e:
        print(f"[ERROR] Drive query failed for: {name} — {e}")
        return None

def download_file(file_id, local_path):
    request = drive_service.files().get_media(fileId=file_id)
    with open(local_path, "wb") as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
    return local_path

def convert_doc_to_docx(doc_path):
    output_path = doc_path + "x"
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    try:
        subprocess.run([soffice_path, "--headless", "--convert-to", "docx", "--outdir", ".", doc_path],
                       check=True, capture_output=True, text=True)
        return output_path if os.path.exists(output_path) else None
    except Exception:
        return None


def extract_text_and_style(docx_path):
    try:
        doc = Document(docx_path)
        song_lines = []

        for p_idx, p in enumerate(doc.paragraphs):
            raw_para = p.text.strip()
            if not raw_para:
                song_lines.append(("", False))
                continue

            # Skip metadata like "Verse 1", "Chorus"
            if re.match(r'^(Hymnal #\d+|Verse \d+|Chorus)$', raw_para.strip(), re.IGNORECASE):
                continue

            for run_idx, run in enumerate(p.runs):
                run_text = run.text.strip()
                if not run_text:
                    continue

                is_italic = run.italic is True

                # Split multiline run content
                for i, line in enumerate(run_text.splitlines()):
                    line = line.strip()
                    if not line:
                        continue
                    song_lines.append((line, is_italic))

        return song_lines

    except Exception as e:
        print(f"[ERROR] extract_text_and_style failed: {e}")
        return []

def update_slide1_right(presentation):
    target_text_prefix = "Corporate Worship Service:"
    upcoming_sunday = datetime.now()
    while upcoming_sunday.weekday() != 6:
        upcoming_sunday += timedelta(days=1)
    date_str = upcoming_sunday.strftime("%B %-d, %Y") if hasattr(upcoming_sunday, 'strftime') else upcoming_sunday.strftime("%B %d, %Y")
    new_text = f"{target_text_prefix} {date_str}"

    # Slide 0 = Slide 1 visually (Python-pptx is 0-indexed)
    slide = presentation.slides[0]

    # Assume right half is shape index 1 or named something predictable
    for shape in slide.shapes:
        if shape.has_text_frame:
            if any(paragraph.text.strip().startswith(target_text_prefix) for paragraph in shape.text_frame.paragraphs):
                for para in shape.text_frame.paragraphs:
                    if para.text.strip().startswith(target_text_prefix):
                        para.text = new_text
                        for run in para.runs:
                            run.font.name = CONFIG["FONT_NAME"]
                            run.font.size = Pt(CONFIG["DATE_FONT_SIZE"])
                        break



def clear_side(slide, side):
    midpoint = Inches(5)
    to_remove = []
    for shape in slide.shapes:
        if not hasattr(shape, "left"):
            continue
        if (side == 'left' and shape.left < midpoint) or (side == 'right' and shape.left >= midpoint):
            to_remove.append(shape)
    for shape in to_remove:
        slide.shapes._spTree.remove(shape._element)
    print(f"[CLEAR] Removed {len(to_remove)} shapes from slide ({side})")

def add_song_content(slide, side, title, lines, slide_index):
    top = Inches(CONFIG["BOX_TOP"])
    width = Inches(CONFIG["BOX_WIDTH"])
    height = Inches(CONFIG["BOX_HEIGHT"])
    left = Inches(CONFIG["LEFT_MARGIN"] if side == 'left' else CONFIG["RIGHT_MARGIN"])

    clear_side(slide, side)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Use tf.paragraphs[0] directly to avoid layout gap
    p_title = tf.paragraphs[0]
    p_title.clear()
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(12)
    p_title.space_before = Pt(0)  # remove extra gap just in case

    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = CONFIG["FONT_NAME"]
    run_title.font.size = Pt(CONFIG["TITLE_SIZE"])
    run_title.font.bold = True

    slide_number = CONFIG["PAGE_NUMBERS"].get((slide_index, side), None)
    if slide_number is not None:
        run_tab = p_title.add_run()
        run_tab.text = "    "
        run_tab.font.size = Pt(CONFIG["TITLE_SIZE"])
        run_tab.font.name = CONFIG["FONT_NAME"]

        run_page = p_title.add_run()
        run_page.text = str(slide_number)
        run_page.font.size = Pt(CONFIG["TITLE_SIZE"])
        run_page.font.name = CONFIG["FONT_NAME"]

    # Remove duplicate title line from lyrics
    lines = lines[1:] if len(lines) > 1 and lines[0][0].strip().lower() == title.strip().lower() else lines

    # Extract footer lines from end
    footer_lines = []
    while lines:
        text = lines[-1][0].strip()
        if not text:
            lines.pop()
            continue
        if re.search(r'(©|CCLI|Public Domain|Words and Music|[0-9]{4}|Translated|by\s+\w+)', text, re.IGNORECASE):
            footer_lines.insert(0, lines.pop())
        else:
            break

    # Main song body
    prev_italic = None  # Track the previous line's italic state

    for (text, is_italic) in lines:
        if not text.strip():
            tf.add_paragraph()
            continue

        # Add spacing when switching from italic to non-italic (e.g. chorus → verse)
        if prev_italic is True and is_italic is False:
            tf.add_paragraph()  # Insert line gap
            print("[SPACING] Inserted blank line between italic and non-italic blocks")

        if '\n' in text:
            for subline in text.splitlines():
                if not subline.strip():
                    tf.add_paragraph()
                    continue
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = subline
                run.font.size = Pt(CONFIG["FONT_SIZE"])
                run.font.name = CONFIG["FONT_NAME"]
                run.font.italic = is_italic
        else:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.size = Pt(CONFIG["FONT_SIZE"])
            run.font.name = CONFIG["FONT_NAME"]
            run.font.italic = is_italic

        prev_italic = is_italic  # Update tracker

    # Footer lines (smaller font)
    for (text, is_italic) in footer_lines:
        p_ref = tf.add_paragraph()
        p_ref.alignment = PP_ALIGN.LEFT
        p_ref.space_before = Pt(6)
        run_ref = p_ref.add_run()
        run_ref.text = text
        run_ref.font.size = Pt(CONFIG["FOOTNOTE_SIZE"])
        run_ref.font.name = CONFIG["FONT_NAME"]
        run_ref.font.italic = is_italic

    print(f"[SUCCESS] Added {title} ({slide_number}) to slide")

def update_order_of_service(prs, song_titles):
    slide = prs.slides[1]  # Slide 2 visually
    left_margin_limit = Inches(5)
    hymn_index = 0

    for shape in slide.shapes:
        if not shape.has_text_frame or shape.left >= left_margin_limit:
            continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text.strip()
            if text.startswith("HYMN") and hymn_index < len(song_titles):
                raw_title = song_titles[hymn_index]
                # Remove any text in brackets for layout safety
                clean_title = re.sub(r'\s*\(.*?\)', '', raw_title).strip()
                quoted_title = f'“{clean_title}”'

                # Estimate padding (adjust total_width if needed)
                total_width = 60
                padding = max(1, total_width - len(quoted_title) - 5)  # 5 for "HYMN"
                hymn_index += 1

                para.clear()

                run_hymn = para.add_run()
                run_hymn.text = "HYMN" + " " * padding
                run_hymn.font.name = CONFIG["FONT_NAME"]
                run_hymn.font.size = Pt(CONFIG["FONT_SIZE"])
                run_hymn.font.bold = True

                run_title = para.add_run()
                run_title.text = quoted_title
                run_title.font.name = CONFIG["FONT_NAME"]
                run_title.font.size = Pt(CONFIG["FONT_SIZE"])
                run_title.font.bold = False

    print(f"[COMPLETE] Inserted {hymn_index} HYMN titles into Order of Service.\n")


def generate_bulletin(song_order):
    prs = Presentation(CONFIG["TEMPLATE_PPTX"])
    update_slide1_right(prs)
    update_order_of_service(prs, song_order)
    folder_id = find_file_id(CONFIG["LYRICS_FOLDER_NAME"])
    if not folder_id:
        print("[ERROR] Lyrics folder not found")
        return

    SONG_SLIDE_MAP = [
        {'slide_index': 3, 'side': 'left', 'song_index': 1},
        {'slide_index': 3, 'side': 'right', 'song_index': 2},
        {'slide_index': 2, 'side': 'left', 'song_index': 3},
        {'slide_index': 1, 'side': 'right', 'song_index': 4},
        {'slide_index': 2, 'side': 'right', 'song_index': 0},
    ]

    temp_files = []
    for entry in SONG_SLIDE_MAP:
        slide = prs.slides[entry['slide_index']]
        side = entry['side']
        idx = entry['song_index']
        if idx >= len(song_order):
            continue
        name = song_order[idx]

        file_id = find_file_id(name + ".docx", parent_id=folder_id) or find_file_id(name + ".doc", parent_id=folder_id)
        if not file_id:
            print(f"[ERROR] File not found for {name}")
            continue

        ext = ".docx" if name + ".docx" in file_id else ".doc"
        local_file = download_file(file_id, name + ext)
        temp_files.append(local_file)

        if ext == ".doc":
            docx_file = convert_doc_to_docx(local_file)
            if docx_file:
                temp_files.append(docx_file)
                local_file = docx_file

        lines = extract_text_and_style(local_file)
        if lines:
            print(f"[INSERT] {name} → Slide {entry['slide_index'] + 1} ({side})")
            add_song_content(slide, side, name, lines, entry['slide_index'])

    for f in temp_files:
        if os.path.exists(f):
            os.remove(f)

    prs.save(CONFIG["OUTPUT_PPTX"])
    print(f"\n✅ Bulletin saved as '{CONFIG['OUTPUT_PPTX']}'")

# -------- RUN --------
if __name__ == "__main__":
    print("Enter songs in order separated by commas (or press Enter to use default):")
    user_input = input().strip()
    order = [s.strip() for s in user_input.split(',')] if user_input else CONFIG["DEFAULT_SONG_ORDER"]
    if not os.path.exists(CONFIG["TEMPLATE_PPTX"]):
        print(f"[ERROR] Template '{CONFIG['TEMPLATE_PPTX']}' not found.")
    else:
        generate_bulletin(order)
