import os
import re
import subprocess
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from docx import Document
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# -------- CONFIG --------
CONFIG = {
    "SERVICE_ACCOUNT_FILE": "credentials.json",
    "LYRICS_FOLDER_NAME": "MPBC Song Lyric Docs",
    "TEMPLATE_PPTX": "BulletinTemplate.pptx",
    "OUTPUT_PPTX": "Updated_Bulletin.pptx",
    "FONT_NAME": "Montserrat",
    "FONT_SIZE": 11,
    "TITLE_SIZE": 11,
    "FOOTNOTE_SIZE": 5,
    "BOX_TOP": 0.1,
    "BOX_WIDTH": 4.8,
    "BOX_HEIGHT": 7,
    "LEFT_MARGIN": 0.1,
    "RIGHT_MARGIN": 5.0,
    "DEFAULT_SONG_ORDER": [
        "Because He Lives",
        "A Mighty Fortress is Our God",
        "Christ the Lord is Risen Today",
        "Christ Arose",
        "See What a Morning (Resurrection Hymn)"
    ]
}

# -------- GOOGLE SETUP --------
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
creds = service_account.Credentials.from_service_account_file(CONFIG["SERVICE_ACCOUNT_FILE"], scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

# -------- HELPERS --------
def find_file_id(name, parent_id=None):
    query = f"name = '{name}'"
    if parent_id:
        query += f" and '{parent_id}' in parents"
    results = drive_service.files().list(q=query, fields="files(id, name)").execute()
    return results.get("files", [])[0]["id"] if results.get("files") else None

def download_file(file_id, local_path):
    request = drive_service.files().get_media(fileId=file_id)
    with open(local_path, "wb") as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
    return local_path

def convert_doc_to_docx(doc_path):
    output_path = doc_path + "x"
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    try:
        subprocess.run([soffice_path, "--headless", "--convert-to", "docx", "--outdir", ".", doc_path],
                       check=True, capture_output=True, text=True)
        return output_path if os.path.exists(output_path) else None
    except Exception:
        return None

def extract_text_and_style(docx_path):
    try:
        doc = Document(docx_path)
        lines = []
        for p in doc.paragraphs:
            paragraph_text = []
            current_italic = False
            for run in p.runs:
                if run.text:
                    paragraph_text.append(run.text)
                    current_italic = current_italic or run.italic
            if paragraph_text:
                full_text = "".join(paragraph_text).strip()
                if full_text and not re.match(r'(Hymnal #\d+|Verse \d+|Chorus)', full_text):
                    lines.append((full_text, current_italic))
        return lines
    except Exception:
        return []

def clear_side(slide, side):
    midpoint = Inches(5)
    to_remove = []
    for shape in slide.shapes:
        if not hasattr(shape, "left"):
            continue
        if (side == 'left' and shape.left < midpoint) or (side == 'right' and shape.left >= midpoint):
            to_remove.append(shape)
    for shape in to_remove:
        slide.shapes._spTree.remove(shape._element)
    print(f"[CLEAR] Removed {len(to_remove)} shapes from slide ({side})")

def add_song_content(slide, side, title, lines):
    top = Inches(CONFIG["BOX_TOP"])
    width = Inches(CONFIG["BOX_WIDTH"])
    height = Inches(CONFIG["BOX_HEIGHT"])
    left = Inches(CONFIG["LEFT_MARGIN"] if side == 'left' else CONFIG["RIGHT_MARGIN"])

    clear_side(slide, side)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    tf.clear()
    p_title = tf.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.size = Pt(CONFIG["TITLE_SIZE"])
    run_title.font.name = CONFIG["FONT_NAME"]
    run_title.font.bold = True
    p_title.alignment = PP_ALIGN.LEFT

    lines = lines[1:] if len(lines) > 1 and lines[0][0].strip().lower() == title.strip().lower() else lines
    last_line_info = lines[-1] if lines and any(c.isdigit() for c in lines[-1][0]) and (',' in lines[-1][0] or '.' in lines[-1][0]) else None
    if last_line_info:
        lines = lines[:-1]

    for (text, is_italic) in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(CONFIG["FONT_SIZE"])
        run.font.name = CONFIG["FONT_NAME"]
        run.font.italic = is_italic

    if last_line_info:
        p_ref = tf.add_paragraph()
        run_ref = p_ref.add_run()
        run_ref.text = last_line_info[0]
        run_ref.font.size = Pt(CONFIG["FOOTNOTE_SIZE"])
        run_ref.font.name = CONFIG["FONT_NAME"]
        run_ref.font.italic = last_line_info[1]
        p_ref.alignment = PP_ALIGN.LEFT
        p_ref.space_before = Pt(12)
        p_ref.space_after = Pt(0)

    print(f"[SUCCESS] Added {title} to slide")

def generate_bulletin(song_order):
    prs = Presentation(CONFIG["TEMPLATE_PPTX"])
    folder_id = find_file_id(CONFIG["LYRICS_FOLDER_NAME"])
    if not folder_id:
        print("[ERROR] Lyrics folder not found")
        return

    SONG_SLIDE_MAP = [
        {'slide_index': 3, 'side': 'left', 'song_index': 1},
        {'slide_index': 3, 'side': 'right', 'song_index': 2},
        {'slide_index': 2, 'side': 'left', 'song_index': 3},
        {'slide_index': 1, 'side': 'right', 'song_index': 4},
        {'slide_index': 2, 'side': 'right', 'song_index': 0},
    ]

    temp_files = []
    for entry in SONG_SLIDE_MAP:
        slide = prs.slides[entry['slide_index']]
        side = entry['side']
        idx = entry['song_index']
        if idx >= len(song_order):
            continue
        name = song_order[idx]

        file_id = find_file_id(name + ".docx", parent_id=folder_id) or find_file_id(name + ".doc", parent_id=folder_id)
        if not file_id:
            print(f"[ERROR] File not found for {name}")
            continue

        ext = ".docx" if name + ".docx" in file_id else ".doc"
        local_file = download_file(file_id, name + ext)
        temp_files.append(local_file)

        if ext == ".doc":
            docx_file = convert_doc_to_docx(local_file)
            if docx_file:
                temp_files.append(docx_file)
                local_file = docx_file

        lines = extract_text_and_style(local_file)
        if lines:
            print(f"[INSERT] {name} → Slide {entry['slide_index'] + 1} ({side})")
            add_song_content(slide, side, name, lines)

    for f in temp_files:
        if os.path.exists(f):
            os.remove(f)

    prs.save(CONFIG["OUTPUT_PPTX"])
    print(f"\n✅ Bulletin saved as '{CONFIG['OUTPUT_PPTX']}'")

# -------- RUN --------
if __name__ == "__main__":
    print("Enter songs in order separated by commas (or press Enter to use default):")
    user_input = input().strip()
    order = [s.strip() for s in user_input.split(',')] if user_input else CONFIG["DEFAULT_SONG_ORDER"]
    if not os.path.exists(CONFIG["TEMPLATE_PPTX"]):
        print(f"[ERROR] Template '{CONFIG['TEMPLATE_PPTX']}' not found.")
    else:
        generate_bulletin(order)
