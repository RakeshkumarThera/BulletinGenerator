import os
import re
import subprocess
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from docx import Document
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# -------- CONFIG --------
SERVICE_ACCOUNT_FILE = 'credentials.json'
LYRICS_FOLDER_NAME = "MPBC Song Lyric Docs"
TEMPLATE_PPTX = "BulletinTemplate.pptx"
OUTPUT_PPTX = "Updated_Bulletin.pptx"
DEFAULT_SONG_ORDER = [
    "Because He Lives",
    "A Mighty Fortress is Our God",
    "Christ the Lord is Risen Today",
    "Christ Arose",
    "See What a Morning (Resurrection Hymn)"
]

# -------- GOOGLE SETUP --------
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
creds = service_account.Credentials.from_service_account_file(SERVICE_ACCOUNT_FILE, scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

# -------- HELPERS --------
def find_file_id(name, parent_id=None):
    query = f"name = '{name}'"
    if parent_id:
        query += f" and '{parent_id}' in parents"
    results = drive_service.files().list(q=query, fields="files(id, name)").execute()
    return results.get("files", [])[0]["id"] if results.get("files") else None

def download_file(file_id, local_path):
    request = drive_service.files().get_media(fileId=file_id)
    with open(local_path, "wb") as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
    return local_path

def convert_doc_to_docx(doc_path):
    output_path = doc_path + "x"
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    try:
        result = subprocess.run([soffice_path, "--headless", "--convert-to", "docx", "--outdir", ".", doc_path],
                                check=True, capture_output=True, text=True)
        print(f"[DEBUG] Conversion stdout: {result.stdout.strip()}")
        print(f"[DEBUG] Conversion stderr: {result.stderr.strip()}")
        return output_path if os.path.exists(output_path) else None
    except subprocess.CalledProcessError as e:
        print(f"[ERROR] Conversion failed for {doc_path}. Command: {' '.join(e.cmd)}. Output: {e.output}. Error: {e.stderr}")
        return None
    except FileNotFoundError:
        print(f"[ERROR] LibreOffice (soffice) not found at {soffice_path}. Please ensure LibreOffice is installed and the path is correct.")
        return None

def extract_text_and_style(docx_path):
    try:
        doc = Document(docx_path)
        lines = []
        for p in doc.paragraphs:
            paragraph_text = []
            current_italic = False
            for run in p.runs:
                text = run.text
                if text:
                    paragraph_text.append(text)
                    current_italic = current_italic or run.italic
            if paragraph_text:
                full_text = "".join(paragraph_text).rstrip()
                if full_text.strip():
                    if not (re.match(r'Hymnal #\d+', full_text) or
                            re.match(r'Verse \d+', full_text) or
                            re.match(r'Chorus', full_text)):
                        lines.append((full_text.strip(), current_italic))
        return lines
    except Exception as e:
        print(f"[ERROR] Failed to extract text from {docx_path}: {e}")
        return []

def clear_side(slide, side, slide_index):
    midpoint = Inches(5)
    shapes_to_remove = [shape for shape in slide.shapes if hasattr(shape, "left") and ((side == 'left' and shape.left.inches < midpoint) or (side == 'right' and shape.left.inches >= midpoint))]
    for shape in shapes_to_remove:
        slide.shapes._spTree.remove(shape._element)
    print(f"[DEBUG] Cleared {len(shapes_to_remove)} shapes on slide {slide_index + 1} ({side})")

def add_song_content(slide, side, title, lines, slide_index):
    if side == 'left':
        left, width = Inches(0.5), Inches(4.5)
        top = Inches(0.5) # Default top
    elif side == 'right':
        left, width = Inches(5.2), Inches(4.5)
        top = Inches(0.5) # Changed from 0.7 to 0.5 to align with left side

    # Remove special offset for Slide 3 (right) since top is now consistent
    # if slide_index == 2 and side == 'right': # Slide 3 is 0-indexed as 2
    #     top = Inches(0.7) # Nudge it down slightly
    #     print(f"[DEBUG] Applying special top offset for Slide {slide_index + 1} ({side}) content box.")

    clear_side(slide, side, slide_index)

    content_box = slide.shapes.add_textbox(left, top, width, Inches(5.7))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    print(f"[DEBUG] Created textbox at left={left.inches}, top={top.inches}, width={width.inches}, height=5.7 for slide {slide_index + 1} ({side})")

    p_title = tf_content.add_paragraph()
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.size = Pt(11)
    run_title.font.name = "Montserrat"
    run_title.font.bold = True
    p_title.alignment = PP_ALIGN.LEFT

    if lines:
        lines_for_content = lines[1:] if len(lines) > 1 and lines[0][0].strip().lower() == title.strip().lower() else lines

        last_line_info = None
        if lines_for_content:
            potential_last_line_text = lines_for_content[-1][0].strip()
            if any(char.isdigit() for char in potential_last_line_text) and (
                    ',' in potential_last_line_text or '.' in potential_last_line_text):
                last_line_info = lines_for_content[-1]
                lines_for_content = lines_for_content[:-1]

        current_verse = []
        for i, (text, is_italic) in enumerate(lines_for_content):
            current_verse.append((text, is_italic))

            should_break_verse = False
            if i < len(lines_for_content) - 1:
                if (text.strip().endswith(('.', '!', '?')) and len(text.strip()) > 0) or \
                   (len(text.split()) <= 2 and not text.strip().endswith(':') and not text.strip().endswith(';')):
                    should_break_verse = True
            elif i == len(lines_for_content) - 1:
                should_break_verse = True

            if should_break_verse:
                p_lyrics = tf_content.add_paragraph()
                p_lyrics.alignment = PP_ALIGN.LEFT
                p_lyrics.space_before = Pt(0)
                p_lyrics.space_after = Pt(0)

                for k, (line_text, line_italic) in enumerate(current_verse):
                    run_lyrics = p_lyrics.add_run()
                    run_lyrics.text = line_text
                    run_lyrics.font.size = Pt(11)
                    run_lyrics.font.name = "Montserrat"
                    run_lyrics.font.italic = line_italic
                    if k < len(current_verse) - 1:
                        run_lyrics.text += '\n'

                if (text.strip().endswith('!') and len(text.strip()) > 0) or \
                   (len(text.split()) <= 2 and not text.strip().endswith(':') and not text.strip().endswith(';')):
                    if i < len(lines_for_content) - 1 or not last_line_info:
                        p_lyrics_space = tf_content.add_paragraph()
                        p_lyrics_space.space_before = Pt(12)
                        p_lyrics_space.space_after = Pt(0)
                current_verse = []

    if last_line_info:
        p_ref = tf_content.add_paragraph()
        run_ref = p_ref.add_run()
        run_ref.text = last_line_info[0]
        run_ref.font.size = Pt(5)
        run_ref.font.name = "Montserrat"
        run_ref.font.italic = last_line_info[1]
        p_ref.alignment = PP_ALIGN.LEFT
        p_ref.space_before = Pt(12)
        p_ref.space_after = Pt(0)

    print(f"[SUCCESS] Added {title} to slide {slide_index + 1} ({side})")


def generate_bulletin(song_order):
    prs = Presentation(TEMPLATE_PPTX)
    lyrics_folder_id = find_file_id(LYRICS_FOLDER_NAME)
    if not lyrics_folder_id:
        print("[ERROR] Lyrics folder not found")
        return

    SONG_SLIDE_MAP = [
        {'slide_index': 3, 'side': 'left', 'song_index': 1},    # A Mighty Fortress is Our God
        {'slide_index': 3, 'side': 'right', 'song_index': 2},   # Christ the Lord is Risen Today
        {'slide_index': 2, 'side': 'left', 'song_index': 3},    # Christ Arose
        {'slide_index': 1, 'side': 'right', 'song_index': 4},   # See What a Morning (Resurrection Hymn)
        {'slide_index': 2, 'side': 'right', 'song_index': 0}    # Because He Lives
    ]

    temp_files = []
    for entry in SONG_SLIDE_MAP:
        try:
            slide = prs.slides[entry['slide_index']]
            print(f"[DEBUG] Processing slide {entry['slide_index'] + 1} ({entry['side']}) with song_index {entry['song_index']}")
        except IndexError:
            print(f"[ERROR] Slide index {entry['slide_index']} is out of range. Check your template.")
            continue

        side = entry['side']
        song_index = entry['song_index']
        if song_index >= len(song_order):
            continue
        song_name = song_order[song_index]
        print(f"[PROCESSING] Slide {entry['slide_index'] + 1} ({side}) → {song_name}")

        file_id = find_file_id(song_name + ".docx", parent_id=lyrics_folder_id)
        local_path = None
        if file_id:
            local_path = download_file(file_id, f"{song_name}.docx")
            temp_files.append(local_path)
        else:
            file_id = find_file_id(song_name + ".doc", parent_id=lyrics_folder_id)
            if file_id:
                doc_path = download_file(file_id, f"{song_name}.doc")
                temp_files.append(doc_path)
                local_path = convert_doc_to_docx(doc_path)
                if local_path:
                    temp_files.append(local_path)
                else:
                    temp_files.remove(doc_path)
        if not local_path:
            print(f"[ERROR] No valid file for {song_name}")
            continue

        lines = extract_text_and_style(local_path)
        if not lines:
            print(f"[DEBUG] No lines extracted for '{song_name}' from '{os.path.basename(local_path)}'. This could be why it's blank.")
        else:
            print(f"[DEBUG] Extracted lines for '{song_name}' from '{os.path.basename(local_path)}': {lines}")

        if lines:
            clear_side(slide, side, entry['slide_index'])
            add_song_content(slide, side, song_name, lines, entry['slide_index'])
        else:
            print(f"[WARNING] No content to add for {song_name} on slide {entry['slide_index'] + 1} ({side})")

    for file in temp_files:
        if os.path.exists(file):
            os.remove(file)
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ Updated bulletin saved as '{OUTPUT_PPTX}'")

# -------- RUN --------
if __name__ == "__main__":
    print("Enter songs in order separated by commas (or press Enter to use default):")
    input_songs = input().strip()
    if input_songs:
        SONG_ORDER = [s.strip() for s in input_songs.split(',') if s.strip()]
    else:
        SONG_ORDER = DEFAULT_SONG_ORDER

    if not os.path.exists(TEMPLATE_PPTX):
        print(f"[ERROR] Template PPTX '{TEMPLATE_PPTX}' not found. Please ensure it's in the same directory as the script.")
    else:
        generate_bulletin(SONG_ORDER)