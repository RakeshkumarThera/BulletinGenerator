import os
import re
import subprocess
import sys
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from docx import Document
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# -------- CONFIG --------
CONFIG = {
    "SERVICE_ACCOUNT_FILE": "credentials.json",
    "LYRICS_FOLDER_NAME": "MPBC Song Lyric Docs",
    "TEMPLATE_PPTX": "BulletinTemplate.pptx",
    "OUTPUT_PPTX": "Updated_Bulletin.pptx",
    "FONT_NAME": "Montserrat",
    "FONT_SIZE": 12,
    "DATE_FONT_SIZE": 10,
    "TITLE_SIZE": 12,
    "FOOTNOTE_SIZE": 5,
    "BOX_TOP": 0.1,
    "BOX_WIDTH": 4.8,
    "BOX_HEIGHT": 7,
    "LEFT_MARGIN": 0.1,
    "RIGHT_MARGIN": 5.0,
    "DEFAULT_SONG_ORDER": [
        "Because He Lives",
        "A Mighty Fortress is Our God",
        "Christ the Lord is Risen Today",
        "Christ Arose",
        "See What a Morning (Resurrection Hymn)"
    ],
    "PAGE_NUMBERS": {
        (1, 'right'): 7,
        (2, 'left'): 6,
        (2, 'right'): 3,
        (3, 'left'): 4,
        (3, 'right'): 5
    }
}

SONG_SLIDE_MAP = [
    {'slide_index': 3, 'side': 'left', 'song_index': 1},
    {'slide_index': 3, 'side': 'right', 'song_index': 2},
    {'slide_index': 2, 'side': 'left', 'song_index': 3},
    {'slide_index': 1, 'side': 'right', 'song_index': 4},
    {'slide_index': 2, 'side': 'right', 'song_index': 0},
]

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
creds = service_account.Credentials.from_service_account_file(CONFIG["SERVICE_ACCOUNT_FILE"], scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)


def find_file_id(name, parent_id=None):
    safe_name = name.replace("'", r"\'")
    query = f"name = '{safe_name}'"
    if parent_id:
        query += f" and '{parent_id}' in parents"
    try:
        print(f"[DEBUG] Querying Drive for: {safe_name}")
        results = drive_service.files().list(q=query, fields="files(id, name)").execute()
        files = results.get('files', [])
        return files[0]['id'] if files else None
    except Exception as e:
        print(f"[ERROR] Drive query failed for: {name} — {e}")
        return None


def validate_song_order(song_order, folder_id):
    missing = []
    for song in song_order:
        docx = find_file_id(song + ".docx", parent_id=folder_id)
        doc = find_file_id(song + ".doc", parent_id=folder_id)
        if not docx and not doc:
            missing.append(song)
    return missing


def read_song_input():
    print("Enter songs in order separated by commas (can be multiline).\nPress Enter twice to finish:")
    lines = []
    while True:
        line = sys.stdin.readline()
        if not line.strip():
            break
        lines.append(line.strip())
    return [s.strip() for s in ','.join(lines).split(',') if s.strip()]


def input_loop(folder_id):
    order = []
    while True:
        if not order:
            order = read_song_input()
            if not order:
                print("[ERROR] No valid songs entered. Please try again.\n")
                continue

        max_required_index = max(entry['song_index'] for entry in SONG_SLIDE_MAP)
        if len(order) <= max_required_index:
            print(f"\n[ERROR] You entered {len(order)} songs, but at least {max_required_index + 1} are required for all slides.")
            print("Please enter more songs.\n")
            order = []
            continue

        used_songs = [order[entry['song_index']] for entry in SONG_SLIDE_MAP]
        missing = validate_song_order(used_songs, folder_id)

        if not missing:
            print("\n[PREVIEW] Slide Assignment:")
            for entry in SONG_SLIDE_MAP:
                print(f" - Slide {entry['slide_index'] + 1} ({entry['side']}): {order[entry['song_index']]}")
            print("")
            return order

        print("\n[ERROR] The following song files were NOT found in Google Drive:")
        for song in missing:
            print(f" - {song}.docx or {song}.doc")

        choice = input("\nWould you like to re-enter just the missing songs? (Y/n): ").strip().lower()
        if choice == 'n':
            order = []
        else:
            for missing_song in missing:
                print(f"Re-enter for '{missing_song}':")
                new_song = input(" → ").strip()
                order = [new_song if s == missing_song else s for s in order]


def download_file(file_id, local_path):
    request = drive_service.files().get_media(fileId=file_id)
    with open(local_path, "wb") as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
    return local_path


def convert_doc_to_docx(doc_path):
    output_path = doc_path + "x"
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    try:
        subprocess.run([soffice_path, "--headless", "--convert-to", "docx", "--outdir", ".", doc_path],
                       check=True, capture_output=True, text=True)
        return output_path if os.path.exists(output_path) else None
    except Exception:
        return None


def extract_text_and_style(docx_path):
    try:
        doc = Document(docx_path)
        song_lines = []
        for p in doc.paragraphs:
            raw_para = p.text.strip()
            if not raw_para:
                song_lines.append(("", False))
                continue
            if re.match(r'^(Hymnal #\d+|Verse \d+|Chorus)$', raw_para.strip(), re.IGNORECASE):
                continue
            for run in p.runs:
                run_text = run.text.strip()
                if not run_text:
                    continue
                is_italic = run.italic is True
                for line in run_text.splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    song_lines.append((line, is_italic))
        return song_lines
    except Exception as e:
        print(f"[ERROR] extract_text_and_style failed: {e}")
        return []

def clear_side(slide, side):
    midpoint = Inches(5)
    to_remove = []
    for shape in slide.shapes:
        if not hasattr(shape, "left"):
            continue
        if (side == 'left' and shape.left < midpoint) or (side == 'right' and shape.left >= midpoint):
            to_remove.append(shape)
    for shape in to_remove:
        slide.shapes._spTree.remove(shape._element)


def add_song_content(slide, side, title, lines, slide_index):
    top = Inches(CONFIG["BOX_TOP"])
    width = Inches(CONFIG["BOX_WIDTH"])
    height = Inches(CONFIG["BOX_HEIGHT"])
    left = Inches(CONFIG["LEFT_MARGIN"] if side == 'left' else CONFIG["RIGHT_MARGIN"])

    clear_side(slide, side)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p_title = tf.paragraphs[0]
    p_title.clear()
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(12)

    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = CONFIG["FONT_NAME"]
    run_title.font.size = Pt(CONFIG["TITLE_SIZE"])
    run_title.font.bold = True

    slide_number = CONFIG["PAGE_NUMBERS"].get((slide_index, side), None)
    if slide_number is not None:
        run_page = p_title.add_run()
        run_page.text = f"    {slide_number}"
        run_page.font.size = Pt(CONFIG["TITLE_SIZE"])
        run_page.font.name = CONFIG["FONT_NAME"]

    lines = lines[1:] if len(lines) > 1 and lines[0][0].strip().lower() == title.strip().lower() else lines

    for (text, is_italic) in lines:
        if not text.strip():
            tf.add_paragraph()
            continue
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(CONFIG["FONT_SIZE"])
        run.font.name = CONFIG["FONT_NAME"]
        run.font.italic = is_italic


def update_slide1_right(presentation):
    target_prefix = "Corporate Worship Service:"
    upcoming_sunday = datetime.now()
    while upcoming_sunday.weekday() != 6:
        upcoming_sunday += timedelta(days=1)
    date_str = upcoming_sunday.strftime("%B %-d, %Y") if hasattr(upcoming_sunday, 'strftime') else upcoming_sunday.strftime("%B %d, %Y")
    new_text = f"{target_prefix} {date_str}"

    slide = presentation.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip().startswith(target_prefix):
                    para.text = new_text
                    for run in para.runs:
                        run.font.name = CONFIG["FONT_NAME"]
                        run.font.size = Pt(CONFIG["DATE_FONT_SIZE"])
                    break


def update_order_of_service(prs, song_titles):
    slide = prs.slides[1]
    left_limit = Inches(5)
    hymn_index = 0
    for shape in slide.shapes:
        if not shape.has_text_frame or shape.left >= left_limit:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text.strip()
            if text.startswith("HYMN") and hymn_index < len(song_titles):
                clean_title = re.sub(r'\s*\(.*?\)', '', song_titles[hymn_index]).strip()
                quoted_title = f'“{clean_title}”'
                total_width = 60
                padding = max(1, total_width - len(quoted_title) - 5)
                hymn_index += 1
                para.clear()
                run_hymn = para.add_run()
                run_hymn.text = "HYMN" + " " * padding
                run_hymn.font.name = CONFIG["FONT_NAME"]
                run_hymn.font.size = Pt(CONFIG["FONT_SIZE"])
                run_hymn.font.bold = True
                run_title = para.add_run()
                run_title.text = quoted_title
                run_title.font.name = CONFIG["FONT_NAME"]
                run_title.font.size = Pt(CONFIG["FONT_SIZE"])


def generate_bulletin(song_order):
    prs = Presentation(CONFIG["TEMPLATE_PPTX"])
    update_slide1_right(prs)
    update_order_of_service(prs, song_order)
    folder_id = find_file_id(CONFIG["LYRICS_FOLDER_NAME"])
    if not folder_id:
        print("[ERROR] Lyrics folder not found")
        return

    temp_files = []
    for entry in SONG_SLIDE_MAP:
        slide = prs.slides[entry['slide_index']]
        side = entry['side']
        idx = entry['song_index']
        if idx >= len(song_order):
            continue
        name = song_order[idx]
        file_id = find_file_id(name + ".docx", parent_id=folder_id) or find_file_id(name + ".doc", parent_id=folder_id)
        if not file_id:
            raise RuntimeError(f"[FATAL] File not found for {name}. This should not happen — input validation must have failed.")
        ext = ".docx" if name + ".docx" in file_id else ".doc"
        local_file = download_file(file_id, name + ext)
        temp_files.append(local_file)
        if ext == ".doc":
            docx_file = convert_doc_to_docx(local_file)
            if docx_file:
                temp_files.append(docx_file)
                local_file = docx_file
        lines = extract_text_and_style(local_file)
        if lines:
            add_song_content(slide, side, name, lines, entry['slide_index'])

    for f in temp_files:
        if os.path.exists(f):
            os.remove(f)

    prs.save(CONFIG["OUTPUT_PPTX"])
    print(f"\n✅ Bulletin saved as '{CONFIG['OUTPUT_PPTX']}'")


if __name__ == "__main__":
    folder_id = find_file_id(CONFIG["LYRICS_FOLDER_NAME"])
    if not folder_id:
        print("[ERROR] Could not find the lyrics folder in Drive.")
        exit()

    final_order = input_loop(folder_id)
    generate_bulletin(final_order)